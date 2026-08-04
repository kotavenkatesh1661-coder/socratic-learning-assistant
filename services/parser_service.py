from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


class DocumentParserError(Exception):
    """Raised when a document cannot be parsed."""


def clean_text(text: str) -> str:
    """
    Clean extracted text by removing extra blank lines
    and unnecessary spaces.
    """
    lines = []

    for line in text.splitlines():
        cleaned_line = " ".join(line.split())

        if cleaned_line:
            lines.append(cleaned_line)

    return "\n".join(lines)


def extract_pdf_text(file_path: str) -> str:
    """Extract text from a PDF file."""
    text_parts = []

    try:
        with fitz.open(file_path) as pdf_document:
            for page in pdf_document:
                page_text = page.get_text("text")

                if page_text:
                    text_parts.append(page_text)

    except Exception as error:
        raise DocumentParserError(
            f"Unable to read the PDF file: {error}"
        ) from error

    return clean_text("\n".join(text_parts))


def extract_docx_text(file_path: str) -> str:
    """Extract text from a Microsoft Word DOCX file."""
    try:
        document = Document(file_path)

        paragraphs = [
            paragraph.text
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]

    except Exception as error:
        raise DocumentParserError(
            f"Unable to read the DOCX file: {error}"
        ) from error

    return clean_text("\n".join(paragraphs))


def extract_pptx_text(file_path: str) -> str:
    """Extract text from a Microsoft PowerPoint PPTX file."""
    text_parts = []

    try:
        presentation = Presentation(file_path)

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(
                    f"Slide {slide_number}\n"
                    + "\n".join(slide_text)
                )

    except Exception as error:
        raise DocumentParserError(
            f"Unable to read the PPTX file: {error}"
        ) from error

    return clean_text("\n\n".join(text_parts))


def extract_txt_text(file_path: str) -> str:
    """Extract text from a plain TXT file."""
    try:
        with open(file_path, "r", encoding="utf-8") as text_file:
            text = text_file.read()

    except UnicodeDecodeError:
        try:
            with open(
                file_path,
                "r",
                encoding="latin-1",
            ) as text_file:
                text = text_file.read()

        except Exception as error:
            raise DocumentParserError(
                f"Unable to decode the TXT file: {error}"
            ) from error

    except Exception as error:
        raise DocumentParserError(
            f"Unable to read the TXT file: {error}"
        ) from error

    return clean_text(text)


def extract_text(file_path: str) -> str:
    """
    Detect the file extension and call the correct parser.
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    parser_functions = {
        ".pdf": extract_pdf_text,
        ".docx": extract_docx_text,
        ".pptx": extract_pptx_text,
        ".txt": extract_txt_text,
    }

    parser_function = parser_functions.get(extension)

    if parser_function is None:
        raise DocumentParserError(
            "Unsupported file type. "
            "Please upload a PDF, DOCX, PPTX, or TXT file."
        )

    extracted_text = parser_function(str(path))

    if not extracted_text.strip():
        raise DocumentParserError(
            "No readable text was found in the uploaded document."
        )

    return extracted_text